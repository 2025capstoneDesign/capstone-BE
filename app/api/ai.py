import subprocess
from fastapi import APIRouter, Depends, File, UploadFile
from fastapi.responses import JSONResponse

from app.model.history import History
from app.service.ai_service import clova_segmentation, download_youtube_audio_segment, generate_note, get_libreoffice_path, analyze_image, analyze_image_kor, process_important_segments, transcribe_audio_file, extract_ppt_text  # 서비스 모듈 임포트
from app.service.mapping_service import LectureSlideMapper, LectureSlideMapperKor
from app.service.auth_service import get_current_user

from app.schema.ai_schema import PptExtractResponse, YouTubeURLRequest, AudioTranscibeResponse
from app.schema.mapping_schema import LectureTextRequest, MappingResultResponse
from app.schema.history_schema import HistoryResponse

from app.model.user import User
from requests import Session
from app.database.session import get_db

from collections import OrderedDict
from pptx import Presentation
from pdf2image import convert_from_path
from fastapi import Form
import tempfile
import os
import shutil
import base64
import io
import logging

logger = logging.getLogger(__name__)
logger.setLevel(logging.DEBUG)

router = APIRouter()

@router.post("/extract-youtube-audio")
async def youtube_audio(request: YouTubeURLRequest):
    """
    YouTube 영상에서 전체 또는 특정 구간의 오디오를 추출
    """
    try:
        if request.start_time and request.end_time:
            # ⏱️ 특정 구간 잘라내기
            from datetime import datetime

            fmt = "%H:%M:%S"
            start = datetime.strptime(request.start_time, fmt)
            end = datetime.strptime(request.end_time, fmt)
            duration_sec = (end - start).total_seconds()
            if duration_sec <= 0:
                raise ValueError("end_time은 start_time보다 뒤여야 합니다.")

            # 초 단위 → HH:MM:SS 포맷
            hours = int(duration_sec // 3600)
            minutes = int((duration_sec % 3600) // 60)
            seconds = int(duration_sec % 60)
            duration_str = f"{hours:02}:{minutes:02}:{seconds:02}"

            file_path = download_youtube_audio_segment(
                url=request.youtube_url,
                start_time=request.start_time,
                duration=duration_str
            )
        else:
            # 🎬 전체 영상에서 오디오 추출
            from app.service.ai_service import download_youtube_audio
            file_path = download_youtube_audio(request.youtube_url)

        return {"message": "YouTube 오디오 다운로드 완료", "data": file_path}
    
    except Exception as e:
        return JSONResponse(
            status_code=500,
            content={"message": f"오디오 추출 실패: {e}", "data": None}
        )


@router.post("/transcribe-audio", response_model=AudioTranscibeResponse)
async def transcribe_audio(audio_file: UploadFile = File(...)):
    """
    업로드된 오디오 파일을 Whisper 모델로 음성을 텍스트로 변환
    """
    try:
        transcribed_text = transcribe_audio_file(audio_file)
        return {"message": "오디오 텍스트 변환 성공", "data": transcribed_text}
    except Exception as e:
        return JSONResponse(status_code=500, content={"message": f"음성 인식 실패: {e}", "data": None})
    

@router.post("/extract-ppt-text", response_model=PptExtractResponse)
def extract_ppt_text_endpoint(ppt_file: UploadFile = File(...)):
    """
    업로드된 PPTX 파일에서 슬라이드별 텍스트 추출 
    """
    try:
        text_list = extract_ppt_text(ppt_file)
        return {"message": "PPT 텍스트 추출 성공", "data": text_list}
    except ValueError as e:
        # 잘못된 입력 등으로 인한 사용자 오류
        return JSONResponse(status_code=400, content={"message": str(e), "data": None})
    except Exception as e:
        return JSONResponse(status_code=500, content={"message": f"PPT 처리 실패: {e}", "data": None})


mapper = LectureSlideMapper()
@router.post("/text-slide-mapping")
def map_lecture_to_slide(request: LectureTextRequest):
    """
    강의 텍스트 + 슬라이드 텍스트를 받아 매핑 후 슬라이드별 세그먼트 리스트로 그룹화
    (segment_index, text, similarity_score 모두 포함)
    """

    # 이미지 캡셔닝 결과를 슬라이드 텍스트로 slide_texts

    # 강의 텍스트를 세그먼트로 분리 (하나의 세그먼트에 10문장 고정)
    segments = mapper.preprocess_and_split_text(request.lecture_text)

    results = mapper.map_lecture_text_to_slides(
        segment_texts=segments, 
        slide_texts=request.slide_texts
    )

    print("강의 녹음본에서 추출한 텍스트 세그먼트 분리 완료")

    # 2. 슬라이드별 세그먼트 리스트로 묶기
    slide_to_segments = {}

    for res in results:
        segment_idx = res["segment_index"]     
        slide_idx = res["matched_slide_index"]  
        similarity_score = res["similarity_score"] 
        slide_key = f"slide{slide_idx}"

        if slide_key not in slide_to_segments:
            slide_to_segments[slide_key] = []

        # 세그먼트 추가
        slide_to_segments[slide_key].append({
            "segment_index": segment_idx,
            "text": segments[segment_idx],
            "similarity_score": round(similarity_score, 4)  # 소수점 4자리로 깔끔하게
        })

    print("슬라이드별 세그먼트 매칭 완료")

    # 3. 최종 결과 반환
    return JSONResponse(content=slide_to_segments)

#################################
mapper_kor = LectureSlideMapperKor()
@router.post("/text-slide-mapping-kor")
def map_lecture_to_slide_kor(request: LectureTextRequest):
    """
    (한국어용) 강의 텍스트 + 슬라이드 텍스트를 받아 매핑 후,
    슬라이드별 세그먼트 리스트로 그룹화 (segment_index, text, similarity_score 포함)
    """

    # 1. 강의 텍스트 세그먼트 분리
    segments = mapper_kor.preprocess_and_split_text_kor(request.lecture_text, 7)

    # 2. 각 세그먼트를 가장 유사한 슬라이드에 매핑
    results = mapper_kor.map_lecture_text_to_slides_kor(
        segment_texts=segments,
        slide_texts=request.slide_texts
    )

    print("🇰🇷 한글 강의 세그먼트 분리 및 매핑 완료")

    # 3. 슬라이드별로 세그먼트 묶기
    slide_to_segments = {}
    for res in results:
        segment_idx = res["segment_index"]
        slide_idx = res["matched_slide_index"]
        similarity_score = res["similarity_score"]
        slide_key = f"slide{slide_idx}"

        if slide_key not in slide_to_segments:
            slide_to_segments[slide_key] = []

        slide_to_segments[slide_key].append({
            "segment_index": segment_idx,
            "text": segments[segment_idx],
            "similarity_score": round(similarity_score, 4)
        })

    print("🇰🇷 슬라이드별 매칭 결과 완료")
    return JSONResponse(content=slide_to_segments)




@router.post("/image-captioning")
async def image_captioning(file: UploadFile = File(...)):
    """
    PPT or PDF -> PDF -> Image -> 슬라이드별 설명 출력 
    """
    try:
        temp_dir = tempfile.mkdtemp()
        uploaded_path = os.path.join(temp_dir, file.filename)

        # 1. 파일 저장
        with open(uploaded_path, "wb") as f:
            shutil.copyfileobj(file.file, f)

        print("파일 저장 완료")

        file_ext = os.path.splitext(file.filename)[-1].lower()

        # 2. 슬라이드 수 확인 (PPTX일 때만)
        if file_ext == ".pptx":
            prs = Presentation(uploaded_path)
            total_slides = len(prs.slides)
        else:
            total_slides = None  # PDF는 슬라이드 수를 추정할 수 없으므로 None

        # 3. PDF 경로 결정
        if file_ext == ".pptx":
            # PPTX -> PDF 변환
            subprocess.run(
                [LIBREOFFICE_PATH, "--headless", "--convert-to", "pdf", "--outdir", temp_dir, uploaded_path],
                check=True
            )
            pdf_path = os.path.join(temp_dir, os.path.splitext(file.filename)[0] + ".pdf")
            if not os.path.exists(pdf_path):
                raise Exception("PDF 변환 실패")
        elif file_ext == ".pdf":
            pdf_path = uploaded_path
        else:
            raise ValueError("지원하지 않는 파일 형식입니다. PDF 또는 PPTX만 허용됩니다.")

        print("PDF 준비 완료")

        # 4. PDF -> 이미지 변환
        images = convert_from_path(
            pdf_path,
            poppler_path=os.getenv('POPPLER_PATH')
        )
        if not images:
            raise Exception("PDF를 이미지로 변환하는 데 실패했습니다.")

        print("PDF -> 이미지 변환 완료")

        print("이미지 캡셔닝 시작")

        results = {}
        if total_slides:
            results["total_slide"] = total_slides

        # 5. 슬라이드 이미지 캡셔닝
        for idx, img in enumerate(images, start=1):
            buffered = io.BytesIO()
            img.save(buffered, format="JPEG")
            img_base64 = base64.b64encode(buffered.getvalue()).decode()
            image_url = f"data:image/jpeg;base64,{img_base64}"

            caption = await analyze_image(image_url)
            results[f"slide{idx}"] = caption

        print("이미지 캡셔닝 완료")
        shutil.rmtree(temp_dir)

        return JSONResponse(content=results)

    except Exception as e:
        logger.debug(e, stack_info=True)
        if 'temp_dir' in locals() and os.path.exists(temp_dir):
            shutil.rmtree(temp_dir)
        return JSONResponse(status_code=500, content={"message": f"처리 실패: {str(e)}"})
    
#############################################

@router.post("/image-captioning-kor")
async def image_captioning(file: UploadFile = File(...)):
    """
    PPT or PDF -> PDF -> Image -> 슬라이드별 설명 출력 
    """
    try:
        temp_dir = tempfile.mkdtemp()
        uploaded_path = os.path.join(temp_dir, file.filename)

        # 1. 파일 저장
        with open(uploaded_path, "wb") as f:
            shutil.copyfileobj(file.file, f)

        print("파일 저장 완료")

        file_ext = os.path.splitext(file.filename)[-1].lower()

        # 2. 슬라이드 수 확인 (PPTX일 때만)
        if file_ext == ".pptx":
            prs = Presentation(uploaded_path)
            total_slides = len(prs.slides)
        else:
            total_slides = None  # PDF는 슬라이드 수를 추정할 수 없으므로 None

        # 3. PDF 경로 결정
        if file_ext == ".pptx":
            # PPTX -> PDF 변환
            subprocess.run(
                [LIBREOFFICE_PATH, "--headless", "--convert-to", "pdf", "--outdir", temp_dir, uploaded_path],
                check=True
            )
            pdf_path = os.path.join(temp_dir, os.path.splitext(file.filename)[0] + ".pdf")
            if not os.path.exists(pdf_path):
                raise Exception("PDF 변환 실패")
        elif file_ext == ".pdf":
            pdf_path = uploaded_path
        else:
            raise ValueError("지원하지 않는 파일 형식입니다. PDF 또는 PPTX만 허용됩니다.")

        print("PDF 준비 완료")

        # 4. PDF -> 이미지 변환
        images = convert_from_path(
            pdf_path,
            poppler_path=os.getenv('POPPLER_PATH')
        )
        if not images:
            raise Exception("PDF를 이미지로 변환하는 데 실패했습니다.")

        print("PDF -> 이미지 변환 완료")

        print("이미지 캡셔닝 시작")

        results = {}
        if total_slides:
            results["total_slide"] = total_slides

        # 5. 슬라이드 이미지 캡셔닝
        for idx, img in enumerate(images, start=1):
            buffered = io.BytesIO()
            img.save(buffered, format="JPEG")
            img_base64 = base64.b64encode(buffered.getvalue()).decode()
            image_url = f"data:image/jpeg;base64,{img_base64}"

            caption = await analyze_image_kor(image_url)
            results[f"slide{idx}"] = caption

        print("이미지 캡셔닝 완료")
        shutil.rmtree(temp_dir)

        return JSONResponse(content=results)

    except Exception as e:
        logger.debug(e, stack_info=True)
        if 'temp_dir' in locals() and os.path.exists(temp_dir):
            shutil.rmtree(temp_dir)
        return JSONResponse(status_code=500, content={"message": f"처리 실패: {str(e)}"})
    



mapper = LectureSlideMapper()
LIBREOFFICE_PATH = get_libreoffice_path()
@router.post("/process-lecture")
async def process_lecture(
    audio_file: UploadFile = File(...),
    doc_file: UploadFile = File(...),
    skip_transcription: bool = Form(False),
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """
    강의 녹음본 + PPT 파일 입력받아서 슬라이드별 세그먼트 매칭 결과 반환환
    """

    try:
        doc_contents = await doc_file.read()

        ## 강의안 download 폴더에 다운로드
        download_path = os.path.join(("download"), doc_file.filename)
        with open(download_path, "wb") as f:
            f.write(doc_contents)


        ### 1. 오디오 파일 텍스트 변환 (또는 스킵)
        if skip_transcription:
            with open(os.path.join("download", "lecture_text.txt"), "r", encoding="utf-8") as f:
                lecture_text = f.read()
        else:
            lecture_text = transcribe_audio_file(audio_file) 
        
        print("오디오 텍스트 변환 완료")

        # 2. PPT 파일 저장 및 변환
        temp_dir = tempfile.mkdtemp()
        file_path = os.path.join(temp_dir, doc_file.filename)

        with open(file_path , "wb") as f:
            f.write(doc_contents)

        file_ext = os.path.splitext(doc_file.filename)[-1].lower()

        # PPT/PDF 처리 
        if file_ext == ".pptx":
            # LibreOffice를 사용하여 PDF로 변환
            subprocess.run([LIBREOFFICE_PATH, "--headless", "--convert-to", "pdf", "--outdir", temp_dir, file_path], check=True)
            file_path = os.path.join(temp_dir, os.path.splitext(doc_file.filename)[0] + ".pdf")

            if not os.path.exists(file_path):
                raise Exception("PDF 변환 실패")

        elif file_ext != ".pdf":
            raise ValueError("지원되지 않는 파일 형식입니다. PPTX 또는 PDF만 허용됩니다.")
        
        images = convert_from_path(
            file_path ,
            poppler_path=os.getenv('POPPLER_PATH')
        )
        print(f"PDF에서 변환된 슬라이드 이미지 수: {len(images)}")

        print("이미지 캡셔닝 시작")

        ### 3. 슬라이드 이미지별 캡션 추출
        slide_captions = []
        for img in images:
            buffered = io.BytesIO()
            img.save(buffered, format="JPEG")
            img_base64 = base64.b64encode(buffered.getvalue()).decode()
            image_url = f"data:image/jpeg;base64,{img_base64}"

            caption = await analyze_image(image_url)
            slide_captions.append(caption)

        print(f"이미지 캡셔닝 결과: {slide_captions}")

        
        print("이미지 캡셔닝 완료")

        # 4. CLOVA Segmentation 사용해서 강의 녹음본을 세그먼트로 나누기 / 고정된 개수의 문장으로 이루어진 세그먼트 사용
        # segments = clova_segmentation(lecture_text)
        segments = mapper.preprocess_and_split_text(lecture_text) 

        print("Clova API 활용 세그먼트 분리 완료")

        # 5. 세그먼트-슬라이드 매핑
        results = mapper.map_lecture_text_to_slides(
            segment_texts=segments,
            slide_texts=slide_captions
        )

        print("세그먼트 슬라이드 매핑 완료")

        # 중요 세그먼트 추출 
        raw_json = {f"segment{i}": seg for i, seg in enumerate(segments)}
        important_segments_result = process_important_segments(raw_json)

        print("매핑 결과 슬라이드 별로 묶기")

        # 6. 매핑 결과를 슬라이드별로 묶기
        slide_to_segments = {}
        for slide_idx in range(len(slide_captions)):
            slide_key = f"slide{slide_idx + 1}"
            slide_to_segments[slide_key] = []

        for res in results:
            segment_idx = res["segment_index"]
            slide_idx = res["matched_slide_index"]
            similarity_score = res["similarity_score"]
            slide_key = f"slide{slide_idx+1}"

            seg_id = f"segment{segment_idx}"
            is_important = seg_id in important_segments_result

            slide_to_segments[slide_key].append({
                "segment_key": seg_id,
                "text": segments[segment_idx],
                "isImportant": str(is_important).lower(),
                "reason": important_segments_result.get(seg_id, {}).get("reason", ""),
                "linkedConcept": "",
                "pageNumber": ""
            })

        sorted_slide_to_segments = OrderedDict(sorted(slide_to_segments.items(), key=lambda x: int(x[0].replace("slide", ""))))

        print("각 슬라이드별 필기 생성 시작")

        final_notes = {}
        for slide_key, segment_list in sorted_slide_to_segments.items():
            slide_idx = int(slide_key.replace("slide", "")) - 1
            slide_caption = slide_captions[slide_idx]
            matched_segment_texts = [seg["text"] for seg in segment_list]
            note_sections, cost = generate_note(slide_caption, matched_segment_texts)
            note_sections["Segments"] = { 
                seg["segment_key"]: {
                    "text": seg["text"],
                    "isImportant": seg["isImportant"],
                    "reason": seg["reason"],
                    "linkedConcept": seg["linkedConcept"],
                    "pageNumber": seg["pageNumber"]
                } for seg in segment_list
            }
            final_notes[slide_key] = note_sections  

        print("각 슬라이드별 필기 생성 완료")

        # 7. 히스토리 DB에 결과 저장
        new_history = History(
            user_email=current_user.email,  
            filename=doc_file.filename,
            notes_json=final_notes,
        )
        db.add(new_history)
        db.commit()
        db.refresh(new_history)

        shutil.rmtree(temp_dir)
        return JSONResponse(content=final_notes)

    
    except Exception as e:
        if 'temp_dir' in locals() and os.path.exists(temp_dir):
            shutil.rmtree(temp_dir)
        return JSONResponse(status_code=500, content={"message": f"처리 실패: {str(e)}"})
