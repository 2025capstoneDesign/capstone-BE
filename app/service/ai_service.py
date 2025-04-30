from fastapi import Path, UploadFile
from pptx import Presentation     # python-pptx: PPT 파일 파싱 라이브러리
from openai import OpenAI
from dotenv import load_dotenv

import os, uuid, shutil
import yt_dlp                     # YouTube 동영상/오디오 다운로드 라이브러리
import whisper                    # OpenAI Whisper 음성 인식 라이브러리
import platform
import requests
import ast

import re
from typing import Tuple, Dict, List
import asyncio
from functools import partial
from pathlib import Path

# 다운로드 및 업로드 파일 저장 폴더
DOWNLOAD_DIR = "download"
os.makedirs(DOWNLOAD_DIR, exist_ok=True)

def download_youtube_audio(url: str) -> str:
    # 저장 오디오 파일명 
    # audio_file_name = f"{uuid.uuid4().hex}.m4a"
    audio_file_name = f"{url}_audio.wav"
    output_path = os.path.join(DOWNLOAD_DIR, audio_file_name)
    # yt_dlp 옵션 설정: 최고 품질의 오디오만 다운로드 (m4a 형식으로 추출)
    ydl_opts = {
        'outtmpl': output_path,
        'format': 'bestaudio/best',
        # 'postprocessors': [{  # 오디오 추출을 위한 FFmpeg 후처리
        #     'key': 'FFmpegExtractAudio',
        #     'preferredcodec': 'm4a'
        # }]
    }
    with yt_dlp.YoutubeDL(ydl_opts) as ydl:
        error_code = ydl.download([url])
    if error_code != 0:
        # yt_dlp 다운로드 실패 (에러 코드가 0이 아니면 오류)
        raise Exception("YouTube 오디오 다운로드 실패")
    return output_path

def transcribe_audio_file(audio_file: UploadFile) -> str:

    # 업로드된 파일(강의 녹음본)을 서버의 DOWNLOAD_DIR에 저장
    file_path = os.path.join(DOWNLOAD_DIR, "audio.wav")

    with open(file_path, "wb") as buffer:
        shutil.copyfileobj(audio_file.file, buffer)
    
    # Whisper 모델 로드 및 음성 인식 수행 (small 모델 사용)
    print("텍스트 변환 진행중")
    model = whisper.load_model("small")
    print(file_path)
    result = model.transcribe(file_path)
    transcribed_text = result.get("text")
    print("텍스트 변환 완료")

    # 강의 녹음본에서 추출한 텍스트 DOWNLOAD_DIR에 저장
    lecture_text_path = os.path.join(DOWNLOAD_DIR, "lecture_text.txt")
    with open(lecture_text_path, "w", encoding="utf-8") as f:
        f.write(transcribed_text)
    print("강의 녹음본에서 추출한 텍스트 " + lecture_text_path + "에 저장 완료")

    return transcribed_text

def extract_ppt_text(ppt_file: UploadFile) -> dict:
    if not ppt_file.filename.lower().endswith(".pptx"):
        raise ValueError("올바른 PPTX 파일을 업로드하세요.")

    # 업로드된 파일(강의안)을 서버의 DOWNLOAD_DIR에 저장
    file_path = os.path.join(DOWNLOAD_DIR, ppt_file.filename)
    with open(file_path, "wb") as buffer:
        shutil.copyfileobj(ppt_file.file, buffer)

    prs = Presentation(file_path)
    
    slides_text = []

    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    texts.append(text)
        slide_text = "\n".join(texts)
        slides_text.append(slide_text)

    # 💡 슬라이드 번호를 키로 한 딕셔너리로 리턴
    return {
        f"슬라이드 {i+1}": text for i, text in enumerate(slides_text)
    }


# 슬라이드 이미지 분석 
load_dotenv()

def get_libreoffice_path():
    system = platform.system()
    if system == "Windows":
        return r"C:\Program Files\LibreOffice\program\soffice.exe"
    elif system == "Darwin":
        return "/Applications/LibreOffice.app/Contents/MacOS/soffice"
    else:
        return "soffice"

client = OpenAI(
    api_key=os.getenv('OPENAI_API_KEY'),
    base_url="https://api.openai.com/v1"
)

def _analyze_image_sync(image_url):
    response = client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {
                "role": "system",
                "content": """Explain slide content following these guidelines:
1. Present as a professor would during class.
2. Focus on key points, avoid unnecessary details not too long.
3. Use narrative prose.
4. Be concise yet informative."""
            },
            {
                "role": "user",
                "content": [
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": image_url,
                            "detail": "low"
                        }
                    }
                ]
            }
        ],
        max_tokens=1000
    )
    return response.choices[0].message.content

async def analyze_image(image_url):
    try:
        loop = asyncio.get_event_loop()
        result = await loop.run_in_executor(
            None,  # 디폴트 쓰레드풀
            partial(_analyze_image_sync, image_url)
        )
        return result
    except Exception as e:
        return f"오류 발생: {str(e)}"

# Clova 문단 나누기 Api를 활용한 세그먼트 분리 
def clova_segmentation(text: str) -> list:
    CLOVA_API_URL = "https://clovastudio.stream.ntruss.com/testapp/v1/api-tools/segmentation"
    CLOVA_API_KEY = os.getenv('CLOVA_API_KEY')

    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {CLOVA_API_KEY}"
    }

    payload = {
        "text": text,
        "alpha": -100,
        "segCnt": -1,
        "postProcess": True,
        "postProcessMaxSize": 500,   #  한 문단의 최대 글자 수
        "postProcessMinSize": 300    #  한 문단의 최소 글자 수
    }

    response = requests.post(CLOVA_API_URL, headers=headers, json=payload)

    if response.status_code != 200:
        raise Exception(f"CLOVA Segmentation API 호출 실패: {response.text}")

    topic_segments = response.json()['result']['topicSeg']
    segments = [" ".join(segment) for segment in topic_segments]  # 문장 리스트를 하나의 세그먼트로 묶기
    return segments





# OpenAI 클라이언트 초기화
client = OpenAI(
    api_key=os.getenv('OPENAI_API_KEY'),
    base_url="https://api.openai.com/v1"
)

# 토큰 비용
TOKEN_COSTS = {
    "gpt-4o": {
        "input": 2.50,
        "output": 10.00
    }
}
EXCHANGE_RATE = 1468.30

def calculate_cost(usage, model="gpt-4o"):
    input_cost = (usage.prompt_tokens / 1_000_000) * TOKEN_COSTS[model]["input"]
    output_cost = (usage.completion_tokens / 1_000_000) * TOKEN_COSTS[model]["output"]
    return (input_cost + output_cost) * EXCHANGE_RATE

def parse_note_sections(note_text: str) -> Dict[str, str]:
    """GPT 출력 결과를 섹션별로 파싱."""
    sections = {
        "Concise Summary Notes": "",
        "Bullet Point Notes": "",
        "Keyword Notes": "",
        "Chart/Table Summary": ""
    }

    # 정규표현식 패턴
    patterns = {
        "Concise Summary Notes": r"1\. Concise Summary Notes",
        "Bullet Point Notes": r"2\. Bullet Point Notes",
        "Keyword Notes": r"3\. Keyword Notes",
        "Chart/Table Summary": r"4\. Chart/Table Summary"
    }

    matches = {}
    for key, pattern in patterns.items():
        match = re.search(pattern, note_text)
        if match:
            matches[key] = match.start()

    # 위치 기준으로 정렬
    sorted_keys = sorted(matches, key=lambda k: matches[k])

    for i, key in enumerate(sorted_keys):
        start = matches[key]
        end = matches[sorted_keys[i+1]] if i+1 < len(sorted_keys) else len(note_text)
        content = note_text[start:end].strip()
        sections[key] = content

    return sections

def generate_note(slide_caption: str, matched_segments: List[str]) -> Tuple[Dict[str, str], float]:
    """슬라이드 캡션 + 매칭된 세그먼트를 기반으로 GPT로 필기 생성하고 파싱."""
    merged_segments = "\n".join(matched_segments)

    prompt = f"""
You are an expert in creating structured notes based on long user inputs.

The user's input consists of:
- A **slide image caption** that summarizes the slide's main content, and
- A set of **matching lecture segments** explaining details related to that slide.

Slide Image Caption:
\"\"\"
{slide_caption}
\"\"\"

Matched Lecture Segments:
\"\"\"
{merged_segments}
\"\"\"

# Important Writing Rules:

**ABSOLUTELY MUST** use the exact following titles, numbered exactly as shown:
   - "1. Concise Summary Notes"
   - "2. Bullet Point Notes"
   - "3. Keyword Notes"
   - "4. Chart/Table Summary"

1. **Concise Summary Notes**  
- Summarize the combined content into natural sentences within 7–8 lines.

2. **Bullet Point Notes**  
- List the key points clearly and briefly in bullet points.  
- Each point should be one sentence or a short phrase.

3. **Keyword Notes**  
- Extract and list around 10 major keywords, concepts, or important terms.  
- Provide a brief explanation for each keyword.

4. **Chart/Table Summary**  
- Try your best to summarize the content in a **chart or table format** if possible.
- A table is especially helpful when listing concepts, comparing items, or explaining step-by-step processes.  
- Only write "Omitted" if it is clearly impossible to express the content in a structured chart or table.

Important writing guidelines you must follow:
- Respond in English if the user input is in English; respond in Korean if the input is in Korean.
- Make the notes concise and clear so that users can understand quickly.
- Eliminate redundant expressions and maintain a logical flow.
- Clearly separate each style of note-taking in the output.
- If a style is not applicable, do not leave it blank; explicitly write Omitted.
- If there are no matching lecture segments for a slide, generate the notes based as much as possible on the slide's image caption alone.
- Each style must be written only once. Do not repeat or duplicate the same style multiple times.

# Output Format Example:

1. Concise Summary Notes
(Your concise summary here)

2. Bullet Point Notes
(Your bullet points here)

3. Keyword Notes
(Your keywords here)

4. Chart/Table Summary
(Your table here or "Omitted")

Now, generate the notes accordingly.
"""
    try:
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "You are a helpful AI assistant specialized in creating structured lecture notes."},
                {"role": "user", "content": prompt}
            ],
            max_tokens=1500,
            temperature=0.3,
        )

        note_text = response.choices[0].message.content.strip()

        print("GPT 응답 원본:\n", note_text)

        note_sections = parse_note_sections(note_text)
        cost = calculate_cost(response.usage)

        return note_sections, cost

    except Exception as e:
        return {"error": f"오류 발생: {str(e)}"}, 0.0




# 중요 세그먼트 찾기 (기준은 사용자가 설정정)
import ast

def process_important_segments(
        raw_json: dict,
        model: str = "gpt-4o",
        min_size: int = 3000,
        last_length: int = 1000
    ) -> dict:
    segments = [
        {"id": int(k.replace("segment", "")), "content": v}
        for k, v in raw_json.items()
    ]
    segments.sort(key=lambda x: x["id"])

    merged, current_pack, current_len = [], [], 0
    for seg in segments:
        current_pack.append(seg)
        current_len += len(seg["content"])
        if current_len >= min_size:
            merged.append({
                "content": "\n".join(
                    [f"{{id : {s['id']} {s['content']}}}" for s in current_pack]
                )
            })
            current_pack, current_len = [], 0
    if current_pack:
        if current_len < last_length and merged:
            merged[-1]["content"] += "\n" + "\n".join(
                [f"{{id : {s['id']} {s['content']}}}" for s in current_pack]
            )
        else:
            merged.append({
                "content": "\n".join(
                    [f"{{id : {s['id']} {s['content']}}}" for s in current_pack]
                )
            })

    important = []
    for block_idx, block in enumerate(merged):
        rsp = client.chat.completions.create(
            model=model,
            temperature=0.7,
            max_tokens=800,
            messages=[
                {"role": "system", "content": """
Analyze the following lecture content and identify only the truly important parts that the professor emphasized.
Each segment is enclosed in curly braces {} and starts with 'id : number'.

Criteria for importance (must meet at least one):
1. Parts explicitly mentioned as being on the exam or test
2. Parts where a core concept is explained in detail AND highlighted as significant
3. Parts the professor repeatedly or strongly emphasized as important
4. Parts prefaced with phrases like "remember this", "this is crucial", or "key point"
5. Parts where the professor spends significantly more time explaining a single concept compared to others (at least 3x longer than average segment length), particularly if accompanied by multiple examples, analogies, or applications

Be selective - not everything is important. Only include segments with clear evidence of importance.
If uncertain about a segment's importance, exclude it.

If you find important parts, respond in List of dictionaries: [{id : number, reason : reason for importance}, ..]
The reason should be concise (max 2-3 sentences) and include specific evidence from the text.
If no important parts are found, return an empty list []
"""},
                {"role": "user", "content": """
{id : 5 (explains a key concept in detail and professor states "this is particularly important and will be on the exam")}
{id : 6 (mentioned as being on the exam with "make sure you understand this for the test")}
{id : 7 (just a brief introduction with no emphasis)}
"""},
                {"role": "assistant", "content": """
[{"id" : 5, "reason" : "Key concept marked as 'particularly important' and confirmed for exam"},
{"id" : 6, "reason" : "Directly mentioned as test material"}]
"""},
                {"role": "user", "content": block["content"]}
            ]
        )
        txt = rsp.choices[0].message.content.strip()
        print(f"\n🟡 [GPT 응답 block {block_idx}]:\n{txt}\n")  # ✅ GPT 응답 출력

        if txt.startswith("[") and txt.endswith("]"):
            try:
                items = ast.literal_eval(txt)
                important.extend(items)
            except Exception as e:
                print(f"❌ 파싱 실패: {e}")  # ✅ 예외 원인 출력
        else:
            print(f"❌ GPT 응답이 [ ] 형식이 아님: {txt[:100]}...")  # 짤막히 미리보기 출력

    return {
        f"segment{item['id']}": {"reason": item["reason"]}
        for item in important
    }
